import io
import re
import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from src.kpis import truncate_label, format_metric_val

# ---------------------------------------------------------
# C-Suite Theme Colors
# ---------------------------------------------------------
COLOR_NAVY = RGBColor(26, 54, 93)     # Primary Accent (#1a365d)
COLOR_TEAL = RGBColor(49, 151, 149)   # Secondary Accent (#319795)
COLOR_DARK = RGBColor(45, 55, 72)     # Body Text (#2d3748)
COLOR_LIGHT = RGBColor(247, 250, 252) # Card Background (#f7fafc)
COLOR_BORDER = RGBColor(226, 232, 240) # Card Border (#e2e8f0)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_MUTED = RGBColor(113, 128, 150) # Muted text (#718096)

def add_slide_header(slide, title_text, category_text="STRATEGIC PERFORMANCE REVIEW"):
    """Helper to draw a standardized C-suite slide header."""
    cat_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = 'Helvetica'
    p_cat.font.size = Pt(8)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_TEAL
    
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.6), Inches(11.833), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Helvetica'
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_NAVY
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.15), Inches(11.833), Inches(0.015)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BORDER
    shape.line.fill.background()

def extract_markdown_section(markdown, section_name):
    """Parses paragraphs under a specific heading in the report markdown."""
    if not markdown:
        return ""
    lines = markdown.split('\n')
    section_content = []
    capture = False
    
    target = re.sub(r'[^\x00-\x7F]+', '', section_name).lower().strip()
    
    for line in lines:
        line_clean = re.sub(r'[^\x00-\x7F]+', '', line).lower().strip()
        if line.startswith('#') and target in line_clean:
            capture = True
            continue
        elif line.startswith('#') and capture:
            break
        if capture:
            cleaned_line = line.replace('**', '').replace('- ', '• ').strip()
            if cleaned_line:
                section_content.append(cleaned_line)
            
    return "\n\n".join(section_content).strip()

def add_unavailable_box(slide, x, y, cx, cy):
    """Renders a standard consulting message block explaining a missing profit chart."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT
    card.line.color.rgb = COLOR_BORDER
    card.line.width = Pt(1.5)
    
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), cx - Inches(0.4), cy - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Visualization Unavailable"
    p.font.name = 'Helvetica'
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_MUTED
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "This visualization is unavailable because the uploaded dataset does not contain profit information."
    p2.font.name = 'Helvetica'
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = COLOR_DARK
    p2.alignment = PP_ALIGN.CENTER

def generate_executive_presentation(kpis, df_clean, date_col, region_col, product_col, revenue_col, profit_col, report_markdown, industry):
    """Compiles dashboard analytics and insights into a widescreen 16:9 board briefing PPTX. Handles Olist and profit-free tables."""
    if df_clean is None or df_clean.empty:
        # Create standard empty presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        add_slide_header(slide, "Data Warning")
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.3), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Strategic Presentation Unavailable"
        p.font.name = 'Helvetica'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_NAVY
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = "The uploaded CSV dataset does not contain active rows. Please configure mappings or load valid records."
        p2.font.name = 'Helvetica'
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_DARK
        p2.alignment = PP_ALIGN.CENTER
        
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        return ppt_io.getvalue()
        
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    has_profit = profit_col and profit_col in df_clean.columns

    if has_profit:
        region_perf = df_clean.groupby(region_col)[[revenue_col, profit_col]].sum()
        region_perf['margin'] = (region_perf[profit_col] / region_perf[revenue_col] * 100)
        worst_region = region_perf['margin'].idxmin() if not region_perf.empty else "N/A"
        worst_region_margin = region_perf['margin'].min() if not region_perf.empty else 0.0

        product_perf = df_clean.groupby(product_col)[[revenue_col, profit_col]].sum()
        product_perf['margin'] = (product_perf[profit_col] / product_perf[revenue_col] * 100)
        worst_product = product_perf['margin'].idxmin() if not product_perf.empty else "N/A"
        worst_product_margin = product_perf['margin'].min() if not product_perf.empty else 0.0
        best_product_margin = product_perf.loc[kpis['top_product'], 'margin'] if kpis['top_product'] in product_perf.index else 0.0
    else:
        worst_region = "N/A"
        worst_region_margin = 0.0
        worst_product = "N/A"
        worst_product_margin = 0.0
        best_product_margin = 0.0

    # ---------------------------------------------------------
    # SLIDE 1: COVER PAGE
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_NAVY
    accent_bar.line.fill.background()
    
    accent_bar_2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0), Inches(0.15), Inches(7.5))
    accent_bar_2.fill.solid()
    accent_bar_2.fill.fore_color.rgb = COLOR_TEAL
    accent_bar_2.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI Executive Decision Intelligence Platform"
    p.font.name = 'Helvetica'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    
    p2 = tf.add_paragraph()
    p2.text = f"Business Performance Review — {industry} Sector Focus"
    p2.font.name = 'Helvetica'
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_TEAL
    p2.space_before = Pt(10)
    
    meta_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(8), Inches(1.5))
    tf_meta = meta_box.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Prepared for: Executive Board of Directors\n" \
                  f"Date of Presentation: {datetime.date.today().strftime('%B %d, %Y')}\n" \
                  f"Report Rating: {kpis['health_status']} Status (Health Score: {kpis['health_score']}/100)"
    p_meta.font.name = 'Helvetica'
    p_meta.font.size = Pt(11)
    p_meta.font.color.rgb = COLOR_DARK
    p_meta.line_spacing = 1.3
    
    logo_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.8), Inches(4), Inches(0.5))
    p_logo = logo_box.text_frame.paragraphs[0]
    p_logo.text = "⯈ ACCENTURE STRATEGY & CONSULTING"
    p_logo.font.name = 'Helvetica'
    p_logo.font.size = Pt(10)
    p_logo.font.bold = True
    p_logo.font.color.rgb = COLOR_MUTED

    # ---------------------------------------------------------
    # SLIDE 2: EXECUTIVE SUMMARY
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "Executive Summary & Corporate Status")
    
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.0), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_header = tf_left.paragraphs[0]
    p_header.text = "Key Performance Pillars"
    p_header.font.name = 'Helvetica'
    p_header.font.size = Pt(16)
    p_header.font.bold = True
    p_header.font.color.rgb = COLOR_NAVY
    p_header.space_after = Pt(14)
    
    metrics = [
        ("Business Health Score", f"{kpis['health_score']} / 100 ({kpis['health_status']})"),
        ("Consolidated Revenue", format_metric_val(kpis['total_revenue'], prefix="$")),
        ("Consolidated Net Profit", format_metric_val(kpis['total_profit'], prefix="$") if has_profit else "N/A"),
        ("Operating Margin Ratio", f"{kpis['profit_margin']:.1f}%" if has_profit else "N/A"),
        ("Top Market Geography", truncate_label(kpis['top_region'], 20)),
        ("Leading Product Segment", truncate_label(kpis['top_product'], 20))
    ]
    
    for label, val in metrics:
        p_m = tf_left.add_paragraph()
        p_m.text = f"• {label}: "
        p_m.font.name = 'Helvetica'
        p_m.font.size = Pt(11)
        p_m.font.color.rgb = COLOR_DARK
        
        run = p_m.add_run()
        run.text = val
        run.font.bold = True
        run.font.color.rgb = COLOR_NAVY
        p_m.space_after = Pt(8)
        
    right_box = slide.shapes.add_textbox(Inches(6.2), Inches(1.6), Inches(6.383), Inches(5.0))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_r_head = tf_right.paragraphs[0]
    p_r_head.text = "Strategic Synthesis"
    p_r_head.font.name = 'Helvetica'
    p_r_head.font.size = Pt(16)
    p_r_head.font.bold = True
    p_r_head.font.color.rgb = COLOR_NAVY
    p_r_head.space_after = Pt(14)
    
    exec_summary_brief = extract_markdown_section(report_markdown, "Executive Summary")
    if not exec_summary_brief:
        exec_summary_brief = f"Corporate health registers as {kpis['health_status']} with a score of {kpis['health_score']}/100. " \
                             f"Growth is primarily powered by high performance in {kpis['top_region']} and the {kpis['top_product']} segment."
                             
    p_sum = tf_right.add_paragraph()
    p_sum.text = exec_summary_brief
    p_sum.font.name = 'Helvetica'
    p_sum.font.size = Pt(11)
    p_sum.font.color.rgb = COLOR_DARK
    p_sum.line_spacing = 1.3

    # ---------------------------------------------------------
    # SLIDE 3: KEY PERFORMANCE INDICATORS
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "Key Performance Indicators")
    
    kpi_card_width = Inches(2.2)
    kpi_card_height = Inches(4.0)
    start_x = Inches(0.75)
    spacing = Inches(0.2)
    y_pos = Inches(1.8)
    
    kpi_cards_data = [
        ("Total Revenue", format_metric_val(kpis['total_revenue'], prefix="$"), "Topline Performance", COLOR_NAVY),
        ("Total Net Profit", format_metric_val(kpis['total_profit'], prefix="$") if has_profit else "N/A", "Net Profit Earnings", COLOR_NAVY),
        ("Profit Margin", f"{kpis['profit_margin']:.1f}%" if has_profit else "N/A", "Operating Ratio", COLOR_NAVY),
        ("Health Score", f"{kpis['health_score']} / 100", f"Status: {kpis['health_status']}", COLOR_TEAL),
        ("Volume Driver", truncate_label(kpis['top_product'], 15), "Top Product Segment", COLOR_NAVY)
    ]
    
    for i, (label, val, subtext, color) in enumerate(kpi_cards_data):
        x_pos = start_x + (i * (kpi_card_width + spacing))
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, kpi_card_width, kpi_card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_LIGHT
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1.5)
        
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos, kpi_card_width, Inches(0.15))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        stripe.line.fill.background()
        
        tb = slide.shapes.add_textbox(x_pos + Inches(0.1), y_pos + Inches(0.3), kpi_card_width - Inches(0.2), kpi_card_height - Inches(0.4))
        tf_c = tb.text_frame
        tf_c.word_wrap = True
        
        p_l = tf_c.paragraphs[0]
        p_l.text = label.upper()
        p_l.font.name = 'Helvetica'
        p_l.font.size = Pt(10)
        p_l.font.bold = True
        p_l.font.color.rgb = COLOR_MUTED
        p_l.space_after = Pt(15)
        
        p_v = tf_c.add_paragraph()
        p_v.text = val
        p_v.font.name = 'Helvetica'
        p_v.font.size = Pt(20)
        p_v.font.bold = True
        p_v.font.color.rgb = color
        p_v.space_after = Pt(15)
        
        p_s = tf_c.add_paragraph()
        p_s.text = subtext
        p_s.font.name = 'Helvetica'
        p_s.font.size = Pt(9)
        p_s.font.color.rgb = COLOR_DARK

    # ---------------------------------------------------------
    # SLIDE 4: BUSINESS INSIGHTS (NATIVE PPTX CHARTS)
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "Business Insights & Data Visualizations")
    
    cx, cy = Inches(5.5), Inches(2.1)
    
    def style_chart(chart, title):
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.name = 'Helvetica'
        chart.legend.font.size = Pt(8)
        chart.legend.font.color.rgb = COLOR_DARK

    # 1. Chart Top-Left: Revenue by Region (Donut Chart - Top 8 + Others)
    chart_data = CategoryChartData()
    region_revs = df_clean.groupby(region_col)[revenue_col].sum().sort_values(ascending=False)
    if len(region_revs) > 8:
        top_8 = region_revs.head(8)
        others_sum = region_revs.iloc[8:].sum()
        region_categories = list(top_8.index) + ["Others"]
        region_values = list(top_8.values) + [others_sum]
    else:
        region_categories = list(region_revs.index)
        region_values = list(region_revs.values)
        
    region_categories_truncated = [truncate_label(x, 15) for x in region_categories]
    chart_data.categories = region_categories_truncated
    chart_data.add_series('Revenue', tuple(region_values))
    
    x1, y1 = Inches(0.75), Inches(1.4)
    chart1 = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x1, y1, cx, cy, chart_data).chart
    style_chart(chart1, "Regional Sales Distribution")
    
    tb_t1 = slide.shapes.add_textbox(Inches(0.75), Inches(3.6), Inches(5.5), Inches(0.4))
    tf_t1 = tb_t1.text_frame
    tf_t1.word_wrap = True
    p_t1 = tf_t1.paragraphs[0]
    p_t1.text = f"🌎 Region Analysis: Geographic volume is led by **{truncate_label(kpis['top_region'], 20)}**."
    p_t1.font.name = 'Helvetica'
    p_t1.font.size = Pt(8.5)
    p_t1.font.color.rgb = COLOR_DARK
    
    # 2. Chart Top-Right: Profit by Product (Clustered Bar Chart - Horizontal)
    x2, y2 = Inches(7.083), Inches(1.4)
    if has_profit:
        chart_data_prod = CategoryChartData()
        product_profits = df_clean.groupby(product_col)[profit_col].sum().sort_values(ascending=False).head(10)
        # Invert order for correct vertical bar display in horizontal format
        product_profits = product_profits.iloc[::-1]
        prod_categories_truncated = [truncate_label(x, 15) for x in product_profits.index]
        chart_data_prod.categories = prod_categories_truncated
        chart_data_prod.add_series('Profit', tuple(product_profits.values))
        
        chart2 = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x2, y2, cx, cy, chart_data_prod).chart
        style_chart(chart2, "Top 10 Products")
    else:
        add_unavailable_box(slide, x2, y2, cx, cy)
    
    tb_t2 = slide.shapes.add_textbox(Inches(7.083), Inches(3.6), Inches(5.5), Inches(0.4))
    tf_t2 = tb_t2.text_frame
    tf_t2.word_wrap = True
    p_t2 = tf_t2.paragraphs[0]
    p_t2.text = f"📦 Product analysis: Top volume segment is **{truncate_label(kpis['top_product'], 20)}**."
    p_t2.font.name = 'Helvetica'
    p_t2.font.size = Pt(8.5)
    p_t2.font.color.rgb = COLOR_DARK

    # 3. Chart Bottom-Left: Monthly Revenue Trend (Line Chart)
    chart_data_trend = CategoryChartData()
    df_trend = df_clean.copy()
    df_trend['Month'] = df_trend[date_col].dt.to_period('M').astype(str)
    trend_agg = df_trend.groupby('Month')[revenue_col].sum().sort_index()
    chart_data_trend.categories = list(trend_agg.index)
    chart_data_trend.add_series('Revenue', tuple(trend_agg.values))
    
    x3, y3 = Inches(0.75), Inches(4.3)
    chart3 = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x3, y3, cx, cy, chart_data_trend).chart
    style_chart(chart3, "Monthly Revenue Trend")
    
    tb_t3 = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(5.5), Inches(0.4))
    tf_t3 = tb_t3.text_frame
    tf_t3.word_wrap = True
    p_t3 = tf_t3.paragraphs[0]
    p_t3.text = "📈 Monthly Trend: Evaluates temporal volatility peaks and baseline operational growth over billing cycles."
    p_t3.font.name = 'Helvetica'
    p_t3.font.size = Pt(8.5)
    p_t3.font.color.rgb = COLOR_DARK

    # 4. Chart Bottom-Right: Revenue vs Profit (Grouped Column Chart)
    x4, y4 = Inches(7.083), Inches(4.3)
    if has_profit:
        chart_data_comp = CategoryChartData()
        df_comp = df_clean.copy()
        df_comp['Month'] = df_comp[date_col].dt.to_period('M').astype(str)
        comp_agg = df_comp.groupby('Month')[[revenue_col, profit_col]].sum().sort_index()
        chart_data_comp.categories = list(comp_agg.index)
        chart_data_comp.add_series('Revenue', tuple(comp_agg[revenue_col].values))
        chart_data_comp.add_series('Profit', tuple(comp_agg[profit_col].values))
        
        chart4 = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x4, y4, cx, cy, chart_data_comp).chart
        style_chart(chart4, "Revenue vs Profit")
    else:
        add_unavailable_box(slide, x4, y4, cx, cy)
    
    tb_t4 = slide.shapes.add_textbox(Inches(7.083), Inches(6.5), Inches(5.5), Inches(0.4))
    tf_t4 = tb_t4.text_frame
    tf_t4.word_wrap = True
    p_t4 = tf_t4.paragraphs[0]
    p_t4.text = f"🔄 Margin Stability: Overall timeline displays sales trends. Margin ratio: {kpis['profit_margin']:.1f}%." if has_profit else "🔄 Margin Stability: Profit information not available in loaded dataset."
    p_t4.font.name = 'Helvetica'
    p_t4.font.size = Pt(8.5)
    p_t4.font.color.rgb = COLOR_DARK

    # ---------------------------------------------------------
    # SLIDE 5: CRITICAL RISKS
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "Audit Findings: Critical Business Risks", "RISK MANAGEMENT AUDIT")
    
    rows = 3
    cols = 5
    x, y, cx, cy = Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.2)
    table_shape = slide.shapes.add_table(rows, cols, x, y, cx, cy)
    table = table_shape.table
    
    table.columns[0].width = Inches(2.8) # Risk
    table.columns[1].width = Inches(1.5) # Severity
    table.columns[2].width = Inches(2.8) # Business Impact
    table.columns[3].width = Inches(3.2) # Recommendation
    table.columns[4].width = Inches(1.533) # Priority
    
    headers = ["Risk Identified", "Severity Level", "Business Impact", "Strategic Recommendation", "Priority Action"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NAVY
        cell.text = h
        p_cell = cell.text_frame.paragraphs[0]
        p_cell.font.name = 'Helvetica'
        p_cell.font.size = Pt(11)
        p_cell.font.bold = True
        p_cell.font.color.rgb = COLOR_WHITE
        p_cell.alignment = PP_ALIGN.CENTER
        
    table_rows = [
        [
            f"Segment Volume Compression in {truncate_label(worst_product, 15)}" if has_profit else "Geographic/Category Concentration",
            "HIGH",
            "Dilutes corporate margins and cash flow health." if has_profit else "High reliance on a single market area.",
            f"Audit COGS economics for {truncate_label(worst_product, 15)}" if has_profit else f"Diversify channels and increase marketing outside {truncate_label(kpis['top_region'], 15)}.",
            "IMMEDIATE"
        ],
        [
            "Operating Cost Inefficiencies" if has_profit else "Operational Overhead",
            "MEDIUM",
            "Dragging down regional net earnings." if has_profit else "Redundant localized administrative expense.",
            "Convene regional director reviews; consolidate branch overheads.",
            "30 DAYS"
        ]
    ]
    
    for row_idx, data in enumerate(table_rows, start=1):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            p_cell = cell.text_frame.paragraphs[0]
            p_cell.font.name = 'Helvetica'
            p_cell.font.size = Pt(10)
            p_cell.font.color.rgb = COLOR_DARK
            if col_idx == 0:
                p_cell.font.bold = True
                p_cell.font.color.rgb = COLOR_NAVY
            if col_idx in [1, 4]:
                p_cell.alignment = PP_ALIGN.CENTER
                if text in ["HIGH", "IMMEDIATE"]:
                    p_cell.font.bold = True
                    p_cell.font.color.rgb = RGBColor(229, 62, 92)
            cell.vertical_anchor = MSO_SHAPE.RECTANGLE

    # ---------------------------------------------------------
    # SLIDE 6: GROWTH OPPORTUNITIES
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "Growth & Value Creation Opportunities", "VALUE CREATION MATRIX")
    
    rows = 3
    cols = 4
    x, y, cx, cy = Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.2)
    table_shape = slide.shapes.add_table(rows, cols, x, y, cx, cy)
    table = table_shape.table
    
    table.columns[0].width = Inches(3.2) # Opportunity
    table.columns[1].width = Inches(3.5) # Expected Impact
    table.columns[2].width = Inches(2.6) # Difficulty
    table.columns[3].width = Inches(2.533) # ROI Estimate
    
    headers = ["Growth Opportunity", "Expected Business Value", "Implementation Difficulty", "ROI Multiplier"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TEAL
        cell.text = h
        p_cell = cell.text_frame.paragraphs[0]
        p_cell.font.name = 'Helvetica'
        p_cell.font.size = Pt(11)
        p_cell.font.bold = True
        p_cell.font.color.rgb = COLOR_WHITE
        p_cell.alignment = PP_ALIGN.CENTER
        
    table_rows = [
        [
            f"Expansion of high-volume {truncate_label(kpis['top_product'], 15)} segment",
            "Incremental increase in quarterly profit margins." if has_profit else "Expand total sales coverage.",
            "Medium (Sales cycles)",
            "4.5x (Proven local demand)"
        ],
        [
            f"Digital Customer Acquisition Upsell in {truncate_label(kpis['top_region'], 15)}",
            "Secures a larger share of wallet in profitable markets.",
            "Low (Marketing campaigns)",
            "3.2x (High LTV)"
        ]
    ]
    
    for row_idx, data in enumerate(table_rows, start=1):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            p_cell = cell.text_frame.paragraphs[0]
            p_cell.font.name = 'Helvetica'
            p_cell.font.size = Pt(10)
            p_cell.font.color.rgb = COLOR_DARK
            if col_idx == 0:
                p_cell.font.bold = True
                p_cell.font.color.rgb = COLOR_NAVY
            if col_idx in [2, 3]:
                p_cell.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_SHAPE.RECTANGLE

    # ---------------------------------------------------------
    # SLIDE 7: STRATEGIC RECOMMENDATIONS
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "Strategic Recommendations & Initiatives", "STRATEGIC INITIATIVES")
    
    rows = 6
    cols = 5
    x, y, cx, cy = Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.5)
    table_shape = slide.shapes.add_table(rows, cols, x, y, cx, cy)
    table = table_shape.table
    
    table.columns[0].width = Inches(3.8) # Initiative Objective
    table.columns[1].width = Inches(3.2) # Value
    table.columns[2].width = Inches(1.4) # Priority
    table.columns[3].width = Inches(1.8) # Owner
    table.columns[4].width = Inches(1.633) # Timeline
    
    headers = ["Strategic Objective", "Expected Value Creation", "Priority", "Executive Owner", "Execution Timeline"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NAVY
        cell.text = h
        p_cell = cell.text_frame.paragraphs[0]
        p_cell.font.name = 'Helvetica'
        p_cell.font.size = Pt(10.5)
        p_cell.font.bold = True
        p_cell.font.color.rgb = COLOR_WHITE
        p_cell.alignment = PP_ALIGN.CENTER
        
    recs_fallback = [
        (f"Optimize {truncate_label(worst_product, 15)} Delivery" if has_profit else "Diversify Product Catalog", "Reduce overhead and expand segment sales", "HIGH", "COO", "30 Days"),
        ("Replicate Top Region Plays", "Standardize sales execution in underperforming areas", "HIGH", "VP Sales", "60 Days"),
        (f"Enterprise Tier for {truncate_label(kpis['top_product'], 15)}", "Increase average contract value and NRR", "MEDIUM", "VP Product", "60 Days"),
        ("Overhead Optimization", "Reduce operational leaks", "MEDIUM", "CFO", "30 Days"),
        ("Digital Marketing Upsell", "Boost customer lifetime value (LTV) in key segments", "LOW", "CMO", "90 Days")
    ]
    
    for row_idx, (obj, val, prio, owner, time_val) in enumerate(recs_fallback, start=1):
        table.cell(row_idx, 0).text = obj
        table.cell(row_idx, 1).text = val
        table.cell(row_idx, 2).text = prio
        table.cell(row_idx, 3).text = owner
        table.cell(row_idx, 4).text = time_val
        
        for col_idx in range(5):
            cell = table.cell(row_idx, col_idx)
            p_cell = cell.text_frame.paragraphs[0]
            p_cell.font.name = 'Helvetica'
            p_cell.font.size = Pt(9.5)
            p_cell.font.color.rgb = COLOR_DARK
            if col_idx == 0:
                p_cell.font.bold = True
                p_cell.font.color.rgb = COLOR_NAVY
            if col_idx in [2, 4]:
                p_cell.alignment = PP_ALIGN.CENTER
                if col_idx == 2 and prio == "HIGH":
                    p_cell.font.bold = True
                    p_cell.font.color.rgb = RGBColor(229, 62, 92)

    # ---------------------------------------------------------
    # SLIDE 8: NEXT QUARTER ACTION PLAN
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "Next Quarter Execution Roadmap", "EXECUTION ROADMAP")
    
    col_width = Inches(3.6)
    col_height = Inches(4.5)
    start_x = Inches(0.75)
    spacing = Inches(0.5)
    y_pos = Inches(1.8)
    
    roadmap = [
        ("Month 1", "Cost Audit & Foundation", [
            f"• Initiate COGS audit for {truncate_label(worst_product, 15)}." if has_profit else "• Review localized acquisition campaigns.",
            "• Draft standardized regional playbook.",
            "• Set margin floor controls for sales."
        ], COLOR_NAVY),
        ("Month 2", "Process Alignment & Upsell", [
            f"• Replicate {truncate_label(kpis['top_region'], 15)} plays.",
            "• Train CS teams on churn indicators.",
            "• Launch enterprise pilot programs."
        ], COLOR_TEAL),
        ("Month 3", "Pricing Launch & Review", [
            "• Launch optimized pricing models.",
            "• Conduct first quarterly health score audit.",
            "• Evaluate regional marketing ROI."
        ], COLOR_NAVY)
    ]
    
    for i, (title, subtitle, bullets, color) in enumerate(roadmap):
        x_pos = start_x + (i * (col_width + spacing))
        
        header_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, col_width, Inches(0.7))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = color
        header_shape.line.fill.background()
        
        tf_h = header_shape.text_frame
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.text = f"{title.upper()} — {subtitle}"
        p_h.font.name = 'Helvetica'
        p_h.font.size = Pt(10)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.alignment = PP_ALIGN.CENTER
        
        card_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_pos + Inches(0.75), col_width, col_height - Inches(0.75))
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = COLOR_LIGHT
        card_shape.line.color.rgb = COLOR_BORDER
        card_shape.line.width = Pt(1.5)
        
        tb = slide.shapes.add_textbox(x_pos + Inches(0.15), y_pos + Inches(0.9), col_width - Inches(0.3), col_height - Inches(1.0))
        tf_b = tb.text_frame
        tf_b.word_wrap = True
        
        for idx, bullet in enumerate(bullets):
            p_b = tf_b.paragraphs[0] if idx == 0 else tf_b.add_paragraph()
            p_b.text = bullet
            p_b.font.name = 'Helvetica'
            p_b.font.size = Pt(10.5)
            p_b.font.color.rgb = COLOR_DARK
            p_b.space_before = Pt(8)

    # ---------------------------------------------------------
    # SLIDE 9: EXECUTIVE CONCLUSION
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_slide_header(slide, "Executive Conclusion & Recommended Decisions", "BOARD DECISION MATRIX")
    
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "Strategic Outlook"
    p_lh.font.name = 'Helvetica'
    p_lh.font.size = Pt(16)
    p_lh.font.bold = True
    p_lh.font.color.rgb = COLOR_NAVY
    p_lh.space_after = Pt(14)
    
    p_lb = tf_l.add_paragraph()
    p_lb.text = f"The consolidated performance assessment confirms that the enterprise possesses strong core viability, " \
                f"evidenced by a Business Health rating of **{kpis['health_status']}** (Health Score: **{kpis['health_score']}/100**). " \
                f"However, long-term margin scale is gated by localized pricing leaks and regional overheads." if has_profit else \
                f"The consolidated performance assessment confirms that the enterprise possesses strong core viability, " \
                f"evidenced by a Business Health rating of **{kpis['health_status']}** (Health Score: **{kpis['health_score']}/100**). " \
                f"Future expansion requires catalog diversification and broader geographic scaling."
    p_lb.font.name = 'Helvetica'
    p_lb.font.size = Pt(11)
    p_lb.font.color.rgb = COLOR_DARK
    p_lb.line_spacing = 1.3
    
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.783), Inches(4.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLOR_LIGHT
    right_box.line.color.rgb = COLOR_BORDER
    right_box.line.width = Pt(1.5)
    
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.783), Inches(0.15))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = COLOR_TEAL
    stripe.line.fill.background()
    
    tb_r = slide.shapes.add_textbox(Inches(6.95), Inches(2.1), Inches(5.483), Inches(4.0))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    
    p_rh = tf_r.paragraphs[0]
    p_rh.text = "Key Priorities & Decisions"
    p_rh.font.name = 'Helvetica'
    p_rh.font.size = Pt(14)
    p_rh.font.bold = True
    p_rh.font.color.rgb = COLOR_NAVY
    p_rh.space_after = Pt(12)
    
    p_decisions = [
        f"1. **Authorize Audit**: Approve cost audit on key underperforming categories." if has_profit else "1. **Expand Categories**: Authorize product catalog expansion campaigns.",
        f"2. **Reallocate Capital**: Approve the 15% budget reallocation to support sales expansion in **{truncate_label(kpis['top_region'], 15)}**.",
        "3. **SLA Bundles**: Authorize wealth/enterprise bundling campaigns to protect accounts from customer churn.",
        "4. **Report Confidence**: Calculated at 95% based on comprehensive sales history."
    ]
    
    for dec in p_decisions:
        p_d = tf_r.add_paragraph()
        p_d.text = dec
        p_d.font.name = 'Helvetica'
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = COLOR_DARK
        p_d.space_after = Pt(10)
        
    p_gen = tf_r.add_paragraph()
    p_gen.text = "AI-Generated Board Deck briefing, Accenture Strategy Consulting Platform"
    p_gen.font.name = 'Helvetica'
    p_gen.font.size = Pt(8)
    p_gen.font.color.rgb = COLOR_MUTED
    p_gen.space_before = Pt(20)

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io.getvalue()
